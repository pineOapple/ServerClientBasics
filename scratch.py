
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

#Setup of the Template environment
template            = r'C:\Users\Noel\Desktop\PPT\Templates\GSTP.pptx'
prs                 = Presentation(template)
img_extr_ppt        = Presentation(r'C:\Users\Noel\Downloads\laösdjf\GSTP_COUPON001.pptx')

content_slide_layout = prs.slide_layouts[1]

#left = top = width = height = Inches(2.0)
#Defining the Structural layout
elements = ['Content'              ,
            'Panel Measurement'    ,
            'SCA-S EL Pre-Laydown' ,
            'SCA-S PL Pre-Laydown' ,
            'SCA-S EL Pre-Laydown' ,
            'SCA-S Alignment'      ,
            'Pre-Dots'             ,
            'Dots'                 ,
            'Adhesive Gap'         ,
            'SCA-S PL Post-Laydown',
            'SCA-S EL Post-Laydown',
            'SCA-S Visual Post-Laydown']

#Adding content slides

for i in elements:
    content_slide    = prs.slides.add_slide(content_slide_layout)
    shapes           = content_slide.shapes
    title_shape      = s
    hapes.placeholders[0]
    title_shape.text = i

#SinglePictureSlides

active_slide = img_extr_ppt.slides[4]
shapes       = active_slide.shapes

for shape in active_slide.placeholders:
    print('%d %s' % (shape.placeholder_format.idx, shape.name))


prs.save(r'C:\Users\Noel\Desktop\PPT\PPs\GSTP_COUPON001.pptx')